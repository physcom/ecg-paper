"""
Türkçe konferans sunumu üretici.

Üretilen dosya:
  C:/Users/enazarkulov/Documents/Мастер/EKG_Conference_Presentation_TR.pptx

İngilizce orijinal sunumun tam çevirisi (23 slayt) — aynı yerleşim, aynı
geçiş ve giriş animasyonları, aynı şekiller. Şekillerin Türkçe varyantı
yoksa (örn. hibrit ablation şekilleri) İngilizce orijinaller kullanılır.

Yardımcı fonksiyonlar İngilizce üreticiden (`generate_conference_materials`)
yeniden kullanılır; slayt içerikleri burada Türkçe olarak yeniden yazılır.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Emu, Inches, Pt

import generate_conference_materials as en

# helpers
new_presentation = en.new_presentation
_blank_layout = en._blank_layout
add_rect = en.add_rect
add_text = en.add_text
add_bullets = en.add_bullets
add_title_banner = en.add_title_banner
add_footer = en.add_footer
add_image_fit = en.add_image_fit
add_table = en.add_table
add_fade_transition = en.add_fade_transition
add_appear_animations = en.add_appear_animations
PRIMARY = en.PRIMARY
ACCENT = en.ACCENT
TEXT = en.TEXT
MUTED = en.MUTED
GREEN = en.GREEN
BG = en.BG

OUT = en.OUT
PPTX_PATH = OUT / "EKG_Conference_Presentation_TR.pptx"


def _fig(en_path: Path) -> Path:
    """Türkçe varyant varsa onu, yoksa İngilizce orijinali döndür."""
    tr_path = en_path.with_name(en_path.stem + "_TR" + en_path.suffix)
    return tr_path if tr_path.exists() else en_path


FIG_COMPARISON = _fig(en.FIG_COMPARISON)
FIG_5000 = _fig(en.FIG_5000)
FIG_1000 = _fig(en.FIG_1000)
FIG_500 = _fig(en.FIG_500)
FIG_500_W4 = _fig(en.FIG_500_W4)
FIG_GEOMETRY = _fig(en.FIG_GEOMETRY)
FIG_HYBRID_HEADLINE = _fig(en.FIG_HYBRID_HEADLINE)
FIG_HYBRID_INFERENCE = _fig(en.FIG_HYBRID_INFERENCE)


# ---------------------------------------------------------------------------
# Slaytlar
# ---------------------------------------------------------------------------

def slide_title(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, sh, fill=RGBColor(0xFF, 0xFF, 0xFF))
    add_rect(slide, 0, Cm(5), sw, Cm(5), fill=PRIMARY)
    add_text(slide, Cm(1.5), Cm(5.3), sw - Cm(3), Cm(1.2),
             "12 Kanallı EKG Sınıflandırmasında Belirleyici Adım:",
             size=Pt(30), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, Cm(1.5), Cm(6.8), sw - Cm(3), Cm(2.4),
             "Anti-Aliasing'li Altörnekleme ile Chapman–Shaoxing üzerinde\n"
             "Temel 1D-CNN ile %88.43'ten %97.34'e",
             size=Pt(20), color=RGBColor(0xDC, 0xE6, 0xF5))
    add_text(slide, Cm(1.5), Cm(11.5), sw - Cm(3), Cm(1),
             "Elaman Nazarkulov", size=Pt(20), bold=True, color=TEXT)
    add_text(slide, Cm(1.5), Cm(12.5), sw - Cm(3), Cm(1.4),
             "Kırgız–Türk Manas Üniversitesi · Bilgisayar Mühendisliği Bölümü\n"
             "Tez danışmanı: Doç. Dr. Bakıt Şarşembayev",
             size=Pt(14), color=MUTED)
    add_text(slide, Cm(1.5), Cm(17), sw - Cm(3), Cm(1),
             "Bahar 2026  ·  ICML formatında konferans sunumu",
             size=Pt(14), italic=True, color=MUTED)
    add_fade_transition(slide)


def slide_motivation(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(
        slide, "Motivasyon",
        subtitle="Sorun olmaması gereken bir baseline'ın ana bulguya dönüşmesinin nedeni",
    )
    items = [
        "12 kanallı EKG, kardiyak aritmi tanısında altın standarttır; ancak "
        "şehir merkezi dışındaki hastanelerde manuel okuma yavaş ve yetersizdir.",
        "Derin 1D-CNN'ler kardiyolog seviyesinde doğruluk sağlar (Rajpurkar 2017, "
        "Hannun 2019) — ancak yayınlanan hatların neredeyse tamamı ham 500 Hz × 10 s = "
        "5000 örneklik sinyali kullanır.",
        "Önceki baseline'ımız (Güz 2025, Chapman–Shaoxing, 78 sınıf, 1D-CNN) "
        "Test Doğruluğu %88.43 / Makro F1 0.8713'te takılı kalmış; 11 sınıf "
        "F1 < 0.60 düzeyinde çökmüştür.",
        "Tez önerisindeki açık hedef: ≥ %90 doğruluk. Doğrudan iyileştirmeler "
        "(attention, focal loss, etiket birleştirme) yol haritasındaydı — ama "
        "doğrulanmaları yavaştı.",
        "Sessiz gözlem: 5000 örnek, QRS-morfolojisi öğrenmek için gerçekten doğru "
        "alıcı alan uzunluğu mu, yoksa sensör örnekleme hızının bir kalıntısı mı?",
    ]
    tb = add_bullets(slide, Cm(1.5), Cm(3), Cm(30), Cm(14), items)
    add_footer(slide, "Elaman Nazarkulov  ·  KTMÜ", "1 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb.shape_id])


def slide_research_question(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Araştırma Sorusu")
    add_rect(slide, Cm(2), Cm(4.5), Cm(30), Cm(4.5), fill=BG)
    add_text(
        slide, Cm(3), Cm(5.3), Cm(28), Cm(3.5),
        "1D-CNN yalnızca QRS ölçeğinde zamansal örüntülere ihtiyaç duyuyorsa, "
        "girdi uzunluğunu agresif biçimde kesebilir miyiz — ve bu tek değişiklik "
        "düz bir baseline'ı elle tasarlanmış bir hibrit modelin ötesine taşıyabilir mi?",
        size=Pt(22), italic=True, color=PRIMARY, align=PP_ALIGN.CENTER,
    )
    items = [
        "H1:  5000 örnek, CNN'in etkin alıcı alan ihtiyacından büyüktür.",
        "H2:  Anti-aliasing'li altörnekleme (scipy.signal.decimate) tanısal içeriği "
        "korurken zamansal boyutu 5×–10× düşürür.",
        "H3:  Baseline ile tam hibrit model arasındaki doğruluk farkı kısmen "
        "uzunluk-optimizasyonu kaynaklıdır — modelleme kaynaklı değildir.",
    ]
    tb = add_bullets(slide, Cm(2), Cm(10), Cm(30), Cm(8), items, size=Pt(20))
    add_footer(slide, "Araştırma sorusu", "2 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb.shape_id])


def slide_related_work(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "İlgili Çalışmalar (özet)")
    items = [
        "Rajpurkar ve ark. 2017, Hannun ve ark. 2019 — kardiyolog seviyesinde "
        "CNN'ler, 91 bin EKG, hepsi tam örnekleme hızında eğitilmiş.",
        "Oh ve ark. 2018 — CNN + LSTM hibrit, %94.8 doğruluk, standart örnekleme "
        "hızında girdi.",
        "Strodthoff ve ark. 2020 (PTB-XL kıyaslaması) — 100 Hz / 1000 örneklik "
        "girdi zaten makro-AUC 0.925'e ulaşıyor. Tam 500 Hz'in zorunlu olmadığına dair ipucu.",
        "Chen 2021, Xu 2022 — referans-düğüm / referans-nokta veri artırma; nadir "
        "sınıf sinyalleri için yararlı ama giriş uzunluğu çalışması yok.",
        "Iwana & Uchida 2021 — kapsamlı zaman serisi veri artırma incelemesi; "
        "örnekleme hızı seçimi sabit hiper-parametre olarak ele alınmış, tasarım değişkeni değil.",
        "Boşluk: aynı koşullarda büyük bir çoklu-etiket 12 kanallı EKG külliyatında "
        "baseline CNN doğruluğunun giriş uzunluğuna göre sistematik karşılaştırması yok.",
    ]
    tb = add_bullets(slide, Cm(1.5), Cm(3), Cm(30.5), Cm(14), items, size=Pt(18))
    add_footer(slide, "İlgili çalışmalar", "3 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb.shape_id])


def slide_dataset(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Veri Seti")
    rows = [
        ("Veri seti", "Kayıt", "Kanal", "Örnekleme", "Rol"),
        ("Chapman–Shaoxing", "45.152", "12", "500 Hz", "Ana deney (bu sunum)"),
        ("PTB-XL", "21.837", "12", "500 Hz", "Veri setleri arası doğrulama (sonraki faz)"),
        ("MIT-BIH Aritmi", "48", "2", "360 → 500 Hz", "Tarihsel referans"),
    ]
    add_table(slide, Cm(1.5), Cm(3.2), Cm(30.5), Cm(4.2), rows, header=True)

    right_items = [
        "Chapman–Shaoxing'te 78 etiket (kök 'X' için 'ECG: X' kopyaları dahil).",
        "Sınıf dengesizliği: ilk 4 sınıf > 5.000 her biri; 30+ sınıf < 50 orijinal.",
        "SQI < 0.85 kayıtlar düşürüldü (ham verinin %6.7'si). Son kullanılabilir: 62.543.",
        "4.500 örnek/sınıf hedefiyle aşırı örnekleme → 353 bin eğitim havuzu "
        "(katmanlı 68/12/20 eğitim/doğrulama/test).",
    ]
    tb = add_bullets(slide, Cm(1.5), Cm(8.5), Cm(30.5), Cm(9), right_items, size=Pt(18))
    add_footer(slide, "Veri seti", "4 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb.shape_id])


def slide_method(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Yöntem — Ön İşleme Hattı",
                     subtitle="Adım 6 (altörnekleme faktörü) hariç tüm konfigürasyonlarda aynı")
    steps = [
        ("1", "500 Hz'e yeniden örnekleme", "sinc interpolasyonu"),
        ("2", "Butterworth BP 0,5–150 Hz + 50 Hz çentik", "şebeke girişimi ve YF gürültü giderimi"),
        ("3", "HP 0,5 Hz taban kayma giderimi", ""),
        ("4", "Kanal başına Z-skor + ±3σ kırpma", ""),
        ("5", "Sabit 10 s segmentler (12 × 5000)", "katmanlı eğitim/doğrulama/test"),
        ("6", "Anti-aliasing'li altörnekleme (yeni)", "scipy.signal.decimate(x, q, ftype='iir', order=8)"),
        ("7", "Referans-düğüm veri artırma", "yaygın için 3×, nadir sınıflar için 10×"),
    ]
    y = Cm(3)
    shape_ids: list[int] = []
    for i, (num, title, sub) in enumerate(steps):
        yy = y + Cm(i * 1.7)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(1.5), yy, Cm(1.1), Cm(1.1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT if title.endswith("(yeni)") else PRIMARY
        circle.line.fill.background()
        tf = circle.text_frame
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        pr = tf.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
        r = pr.add_run(); r.text = num
        r.font.size = Pt(16); r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tb = add_text(slide, Cm(3.2), yy + Cm(0.1), Cm(28), Cm(1.5),
                      title, size=Pt(18), bold=True,
                      color=ACCENT if title.endswith("(yeni)") else TEXT)
        if sub:
            add_text(slide, Cm(3.2), yy + Cm(0.85), Cm(28), Cm(1.0),
                     sub, size=Pt(13), color=MUTED, italic=True)
        shape_ids.append(tb.shape_id)
    add_footer(slide, "Yöntem — hat", "5 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, shape_ids)


def slide_decimation(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Ana Değişiklik: Anti-Aliasing'li Altörnekleme",
                     subtitle="Ön işlemeye üç satırlık bir yama — başka hiçbir şey değişmiyor.")
    code = (
        "from scipy.signal import decimate\n\n"
        "# faktor = 5  →  5000 → 1000 örnek  (100 Hz etkin)\n"
        "# faktor = 10 →  5000 →  500 örnek  ( 50 Hz etkin)\n"
        "x_down = decimate(x, factor,\n"
        "                  ftype='iir',   # Chebyshev tip-I alt-geçiren\n"
        "                  n=8,           # filtre derecesi\n"
        "                  zero_phase=True)"
    )
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.5), Cm(3.2), Cm(18), Cm(7))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.6); tf.margin_top = Cm(0.4)
    tf.paragraphs[0].text = ""
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.name = "Consolas"; r.font.size = Pt(16)
        r.font.color.rgb = RGBColor(0xE6, 0xE6, 0xEC)

    right = [
        "Fazı korur (zero_phase=True).",
        "Chebyshev tip-I, yeni Nyquist altında anti-aliasing yapar; QRS "
        "kompleksinin yapısal örtüşmesi olmaz.",
        "Filtre çevrimdışı ve bir kez uygulanır — eğitim maliyeti yoktur.",
        "Parametreler ve veri hijyeni 5000 örneklik baseline ile aynıdır "
        "(aynı seed, veri artırma, sınıf ağırlıkları).",
    ]
    tb = add_bullets(slide, Cm(20), Cm(3.5), Cm(12), Cm(10), right, size=Pt(16))
    add_footer(slide, "Yöntem — altörnekleme", "6 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [box.shape_id, tb.shape_id])


def slide_geometry(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(
        slide,
        "Neden İşe Yarıyor — Geometrik Değişmezlik",
        subtitle="scipy.signal.decimate'in koruduğu şey: referans-nokta grafı "
                 "(atım başına P / Q / R / S / T).",
    )
    add_image_fit(
        slide, FIG_GEOMETRY,
        Cm(1.5), Cm(2.7), Cm(30.8), Cm(11.0),
        caption=None,
    )
    cards = [
        ("Alıcı alan",
         "2048 örneklik CNN penceresi artık girdinin %100'ünü kaplıyor; "
         "len=5000'de yaklaşık %40 idi."),
        ("Referans-nokta yoğunluğu",
         "Yaklaşık 60 P/Q/R/S/T noktası len=500 girdisinin %12'sini kaplar; "
         "len=5000'de %1,2 idi — referans nokta başına 10× daha fazla gradyan sinyali."),
        ("Parametre ekonomisi",
         "Aynı 3,72 M parametre, baseline modellemeden LVH, AV-blok, atriyal "
         "flutter gibi sınıfların morfolojisine yeniden tahsis edilir."),
    ]
    box_y = Cm(14.2)
    box_w = Cm(9.8)
    box_h = Cm(4.0)
    xs = [Cm(1.5), Cm(11.7), Cm(21.9)]
    for (title, body), x in zip(cards, xs):
        add_rect(slide, x, box_y, box_w, box_h, fill=BG)
        add_text(
            slide, x + Cm(0.5), box_y + Cm(0.35), box_w - Cm(1.0), Cm(0.9),
            title, size=Pt(14), bold=True, color=PRIMARY,
        )
        add_text(
            slide, x + Cm(0.5), box_y + Cm(1.25), box_w - Cm(1.0), Cm(2.55),
            body, size=Pt(12), color=TEXT,
        )

    add_footer(slide, "Neden işe yarıyor · geometrik değişmezlik", "7 / 21")
    add_fade_transition(slide)


def slide_experimental_setup(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Deneysel Düzenek",
                     subtitle="Dört konfigürasyon yalnızca giriş uzunluğu / DataLoader işçilerinde farklı")
    rows = [
        ("Konfig", "Giriş", "Parametre", "Epoch (s)", "Erken durma", "Log dosyası"),
        ("len=5000", "12 × 5000", "3,72 M", "≈ 195",  "epoch 92", "results-22-04-2026.txt"),
        ("len=1000", "12 × 1000", "3,72 M", "≈ 32",  "epoch 92", "result-23-04-2026-1000.txt"),
        ("len=500",  "12 × 500",  "3,72 M", "≈ 30",  "epoch 100","result-22-04-2026-500.txt"),
        ("len=500 + 4 işçi", "12 × 500", "3,72 M", "≈ 20", "epoch 100",
         "result-23-04-2026-500-4-workers.txt"),
    ]
    add_table(slide, Cm(1.5), Cm(3), Cm(30.5), Cm(5.5), rows, header=True)

    rest = [
        "GPU: NVIDIA RTX 5090, CUDA 12.8, AMP (FP16).",
        "Kayıp: İkili çapraz entropi (çoklu etiket); ters frekans sınıf ağırlıkları.",
        "Optimizer: Adam β1=0,9, β2=0,999, başlangıç LR 1e-3, ReduceLROnPlateau sabır 5.",
        "Erken durma sabrı: doğrulama kaybı üzerinde 10 epoch. Maksimum 100 epoch.",
        "Seed ve veri ayırmaları konfig boyunca sabit — doğrudan karşılaştırma için.",
    ]
    tb = add_bullets(slide, Cm(1.5), Cm(9.5), Cm(30.5), Cm(8), rest, size=Pt(17))
    add_footer(slide, "Deneysel düzenek", "8 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb.shape_id])


def slide_results_table(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Bulgular — Manşet Rakamlar")
    rows = [
        ("Konfig", "Test Doğr.", "Makro P", "Makro R", "Makro F1", "Çıkarım", "Güven"),
        ("len=5000 (baseline)", "%88,43", "0,8768", "0,8847", "0,8713", "89,88 ms", "%12,89"),
        ("len=1000",            "%97,22", "0,9740", "0,9729", "0,9716", "26,14 ms", "%68,88"),
        ("len=500",             "%97,34", "0,9734", "0,9749", "0,9737", "27,20 ms", "%76,23"),
        ("len=500 + 4 işçi",    "%97,38", "0,9741", "0,9751", "0,9744", "43,50 ms", "%69,59"),
    ]
    add_table(slide, Cm(1.5), Cm(3), Cm(30.5), Cm(7), rows)

    add_text(
        slide, Cm(1.5), Cm(11), Cm(30.5), Cm(1.2),
        "Δ (len=5000 → len=500): doğruluk +8,91 pp,  makro F1 +10,24 pp,  çıkarım 3,3× daha hızlı.",
        size=Pt(20), bold=True, color=GREEN,
    )
    add_text(
        slide, Cm(1.5), Cm(12.5), Cm(30.5), Cm(4),
        "Tez önerisinden hedef (≥ %90) ve v2 hedefi (referans-düğüm veri artırma "
        "ile tam attention-CNN-LSTM hibrit için %94,8) — düz baseline'a yalnızca "
        "bir ön işleme değişikliğiyle aşılmıştır.",
        size=Pt(16), italic=True, color=TEXT,
    )
    add_footer(slide, "Bulgular — manşet", "9 / 21")
    add_fade_transition(slide)


def slide_results_figure(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Bulgular — Karşılaştırma Şekli")
    add_image_fit(
        slide, FIG_COMPARISON, Cm(2), Cm(3), Cm(30), Cm(14),
        caption="Şekil 1. Giriş uzunluğuna göre temel 1D-CNN performansı. Dört konfigürasyon "
                "için test doğruluğu, makro F1 ve çıkarım süresi (Chapman–Shaoxing, 78 sınıf).",
    )
    add_footer(slide, "Bulgular — şekil", "10 / 21")
    add_fade_transition(slide)


def slide_training_histories(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Eğitim Dinamiği — Yan Yana 4 Konfig")
    panels = [
        (FIG_5000,  "len=5000"),
        (FIG_1000,  "len=1000"),
        (FIG_500,   "len=500"),
        (FIG_500_W4,"len=500 + 4 işçi"),
    ]
    positions = [
        (Cm(1.0), Cm(2.8)),
        (Cm(17.5), Cm(2.8)),
        (Cm(1.0), Cm(10.5)),
        (Cm(17.5), Cm(10.5)),
    ]
    for (path, label), (x, y) in zip(panels, positions):
        if path.exists():
            pic = slide.shapes.add_picture(str(path), x, y, width=Cm(15.5))
            if pic.height > Cm(7.0):
                ratio = Cm(7.0) / pic.height
                pic.height = Cm(7.0)
                pic.width = int(pic.width * ratio)
        add_text(slide, x, y - Cm(0.6), Cm(15.5), Cm(0.5),
                 label, size=Pt(14), bold=True, color=PRIMARY)
    add_footer(slide, "Eğitim dinamiği", "11 / 21")
    add_fade_transition(slide)


def slide_per_class(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Sınıf Bazında İyileşme",
                     subtitle="Baseline'ın kırıldığı yerde, altörnekleme onu neredeyse tamamen düzeltiyor")
    rows = [
        ("Sınıf (örnek sayısı = 900-1000)", "len=5000 F1", "len=500 F1", "Δ"),
        ("Sol Ventriküler Hipertrofi (LVH)", "0,022", "≥ 0,99", "+0,97"),
        ("EKG: Q dalgası anormalliği",       "0,180", "≥ 0,99", "+0,81"),
        ("İç iletim farkı / IV blok",        "0,286", "≥ 0,98", "+0,70"),
        ("Atriyoventriküler blok",           "0,324", "0,984",  "+0,66"),
        ("Erken atriyal kontraksiyon",       "0,329", "≥ 0,97", "+0,64"),
        ("EKG: atriyal fibrilasyon",         "0,436", "≥ 0,95", "+0,51"),
        ("EKG: ST segment değişiklikleri",   "0,457", "≥ 0,96", "+0,50"),
        ("Birinci derece AV blok",           "0,497", "≥ 0,96", "+0,46"),
        ("EKG: atriyal flutter",             "0,581", "≥ 0,99", "+0,41"),
    ]
    add_table(slide, Cm(1.5), Cm(3), Cm(30.5), Cm(11), rows)
    add_text(
        slide, Cm(1.5), Cm(14.5), Cm(30.5), Cm(2),
        "Tüm 11 hatalı sınıf F1 değeri, modele, kayba veya veri artırma reçetesine "
        "dokunulmadan ≥ 0,95'a yükselir. Darboğaz, girdiydi.",
        size=Pt(16), italic=True, color=TEXT,
    )
    add_footer(slide, "Sınıf bazında iyileşme", "12 / 21")
    add_fade_transition(slide)


# ---------------------------------------------------------------------------
# Hibrit-plan ablation slaytları (30 Nisan 2026)
# ---------------------------------------------------------------------------

def slide_hybrid_intro(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(
        slide,
        "Hibrit-plan ablation: iki başlık taahhüdünün ayrıştırılması",
        subtitle="Başlıktaki iki bağımsız taahhüdün etkisini ölçmek için 2×2 ablation gereklidir.",
    )
    items = [
        "Tez başlığı iki bağımsız taahhüt içerir: (i) 12 kanallı girdi ve "
        "(ii) altörneklenmiş baseline üzerinde referans-düğüm veri artırma.",
        "30 Nisan 2026'ya kadar bu ikisi her zaman birlikte değiştirildi — bu yüzden "
        "doğruluğu hangi kolun gerçekten taşıdığını söyleyemiyorduk.",
        "Bu nedenle 2×2 ablation çalıştırıldı: {1 kanal, 12 kanal} × "
        "{artırma KAPALI, artırma AÇIK}, sabit altörnekleme=10 (len=500), "
        "aynı seed, aynı kod sürümü, Chapman–Shaoxing.",
        "Her hücre, daha önce bildirilen len=500 referansıyla aynıdır: aynı 1D-CNN, "
        "aynı optimizer, aynı eğitim/doğrulama/test ayırması, aynı erken durma ölçütü.",
        "Dört hücre şu soruyu cevaplar: 'başlığın hibrit planı +9 pp kazancın gerçek "
        "kaynağı mı, yoksa içeriklerden biri tüm işi mi yapıyor?'",
    ]
    tb = add_bullets(slide, Cm(1.5), Cm(3), Cm(30.5), Cm(14), items, size=Pt(18))
    add_footer(slide, "Hibrit ablation · düzenek", "13 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb.shape_id])


def slide_hybrid_augment(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(
        slide,
        "Baskın kol: veri artırma",
        subtitle="Referans-düğüm veri artırmasını açıp kapatmak doğruluğu ≈30 pp "
                 "değiştirir — 12 kanal olsun ya da olmasın.",
    )
    add_image_fit(
        slide, FIG_HYBRID_HEADLINE,
        Cm(1.0), Cm(3.0), Cm(17.0), Cm(13.0),
        caption=None,
    )
    items = [
        "artırma KAPALI → AÇIK: 1 kanal +30,36 pp doğruluk / +0,91 makro F1 "
        "(%67,14 → %97,50, 0,0682 → 0,9755).",
        "artırma KAPALI → AÇIK: 12 kanal +29,11 pp doğruluk / +0,90 makro F1 "
        "(%68,29 → %97,40, 0,0762 → 0,9743).",
        "Referans-düğüm yöntemi (oversampler) başlığın ampirik kalbidir — "
        "kanal sayısı değil.",
        "Onsuz, uzun kuyruklu sınıflar makro F1'i 0,08'in üzerine taşımak için "
        "yeterli gradyanı asla biriktiremez — kaç kanal beslersek besleyelim.",
    ]
    tb = add_bullets(slide, Cm(18.5), Cm(3.5), Cm(14), Cm(13), items, size=Pt(16))
    add_footer(slide, "Hibrit ablation · veri artırma kolu", "14 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb.shape_id])


def slide_hybrid_channels(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(
        slide,
        "Kanal sayısı bir kol değil, eşitliği bozan bir faktör",
        subtitle="Veri artırma açıkken, 1 kanal ile 12 kanal seed gürültüsü düzeyinde.",
    )
    rows = [
        ("Konfig", "Test Doğr.", "Makro F1", "Çıkarım", "Güven"),
        ("1 kanal · artırma KAPALI",   "%67,14", "0,0682", "14,74 ms", "%60,01"),
        ("12 kanal · artırma KAPALI",  "%68,29", "0,0762", "14,76 ms", "%87,88"),
        ("1 kanal · artırma AÇIK",     "%97,50", "0,9755", "13,25 ms", "%77,38"),
        ("12 kanal · artırma AÇIK",    "%97,40", "0,9743", "45,86 ms", "%90,15"),
    ]
    add_table(slide, Cm(1.0), Cm(3.0), Cm(16.5), Cm(7.0), rows)

    items = [
        "Artırma AÇIK iken: 1 kanal, 12 kanaldan 0,10 pp doğruluk ve 0,0012 makro F1 "
        "öndedir — bu külliyatın seed gürültü bandı içinde.",
        "Artırma KAPALI iken: 12 kanal, 1 kanaldan 1,15 pp doğruluk ve 0,008 makro F1 "
        "öndedir — yine seed gürültü düzeyinde.",
        "Altörnekleme tanısal sinyali referans-nokta grafında zaten yoğunlaştırmıştır — "
        "model aynı morfolojiyi geri kazanmak için 12 kanala ihtiyaç duymaz.",
        "Veri artırma altörneklenmiş girdiye uygulandığında, kanal sayısı doğruluk "
        "kolu değil, dağıtım seçimidir.",
    ]
    tb = add_bullets(slide, Cm(18.0), Cm(3.0), Cm(14.5), Cm(14), items, size=Pt(15))
    add_footer(slide, "Hibrit ablation · kanal sayısı", "15 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb.shape_id])


def slide_hybrid_deployment(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(
        slide,
        "Dağıtım çıkarımı: kenar (edge) vs hastane",
        subtitle="Aynı doğruluk, farklı işletme noktası: hız vs güven.",
    )
    has_fig = FIG_HYBRID_INFERENCE.exists()
    if has_fig:
        add_image_fit(
            slide, FIG_HYBRID_INFERENCE,
            Cm(1.0), Cm(3.0), Cm(17.0), Cm(13.0),
            caption=None,
        )
        right_x = Cm(18.5)
        right_w = Cm(14)
    else:
        right_x = Cm(1.5)
        right_w = Cm(30.5)
    items = [
        "1 kanal, tek örnek CPU çıkarımında 12 kanaldan 3,5× daha hızlıdır "
        "(13,25 ms vs 45,86 ms).",
        "12 kanal, ayrılmış bir örnekte daha yüksek softmax güveni verir "
        "(%90,15 vs %77,38) — güven-eşiği okumaları için daha geniş marj.",
        "Kenar / giyilebilir dağıtım → 1 kanal: daha hızlı, aynı doğruluk, daha düşük "
        "donanım ve bant genişliği bütçesi.",
        "Karar-başına güven raporlamalı hastane iş akışı → 12 kanal: doğruluk aynı, "
        "ancak kalibrasyon marjı daha geniştir.",
        "Başlığın '12 kanal' taahhüdü bu nedenle tek bir sınıflandırıcı ailesinin "
        "hastane-sınıfı işletme noktası olarak okunmalıdır.",
    ]
    tb = add_bullets(slide, right_x, Cm(3.5), right_w, Cm(13), items, size=Pt(16))
    add_footer(slide, "Hibrit ablation · dağıtım", "16 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb.shape_id])


def slide_speed(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Hız: Eğitim ve Çıkarım",
                     subtitle="Daha kısa girdinin yan etkisi: deneyler etkileşimli hale geliyor")
    left = [
        "Epoch süresi (Chapman–Shaoxing, RTX 5090, AMP):",
        "   len=5000        ≈ 195 s / epoch",
        "   len=1000        ≈  32 s / epoch  (6,1× daha hızlı)",
        "   len=500         ≈  30 s / epoch  (6,5× daha hızlı)",
        "   len=500 + 4 iş. ≈  20 s / epoch  (9,8× daha hızlı)",
        "",
        "Çıkarım (tek 12 kanallı örnek):",
        "   len=5000:  89,88 ms",
        "   len=500:   27,20 ms  (3,3× daha hızlı)",
    ]
    right = [
        "Hiper-parametre ayarı etkileşimli bir döngü oluyor (≈ 10 dk koşular).",
        "DataLoader işçi sayısı=4, farkında olmadığımız bir G/Ç darboğazını ortaya "
        "çıkarıyor — model değişmeden epochda +%33 hızlanma.",
        "Çıkarım, klinik 1 saniyelik bütçenin çok altında, CPU sınıfı marjda.",
        "Model dosyası boyut ve parametre sayısı olarak değişmez; CPU / kenar "
        "dağıtım da ≈3× çıkarım hızlanmasını bedavaya alır.",
    ]
    tb1 = add_bullets(slide, Cm(1.5), Cm(3), Cm(15.5), Cm(14), left,
                      size=Pt(17), bullet="")
    tb2 = add_bullets(slide, Cm(17), Cm(3), Cm(15.5), Cm(14), right, size=Pt(16))
    add_footer(slide, "Hız", "17 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb1.shape_id, tb2.shape_id])


def slide_discussion(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Neden Bu İşe Yarıyor?")
    items = [
        "Alıcı alan. 1D-CNN'imizin son katmanı len=5000'de ≈ 2048 giriş örneğini "
        "görür — 10 s'lik pencerenin büyük kısmını. QRS özellikleri rahatça ≈ 200 ms "
        "içine sığar (500 Hz'te 100 örnek). Daha kısa girdi gradyan sinyalini "
        "fizyolojik olarak ilgili örüntülere yoğunlaştırır.",
        "Parametre ekonomisi. Aynı 3,72 M parametre artık 10× daha az zaman adımını "
        "modelliyor. Kapasite sınıf ayrımları için harcanır, gereksiz zamansal "
        "doldurmaya değil.",
        "Etiket kopyaları kısmen soğurulur. 'EKG: atriyal flutter' ve 'Atriyal flutter' "
        "len=500'de F1 ≥ 0,95'a ulaşır — model uzun, ilgisiz uzantıları ezberlemek "
        "zorunda kalmadığında her ikisini de ezberleyebilir.",
        "Anti-aliasing önemlidir. 10'a basit strided pooling katlanmış spektrum üretir; "
        "Chebyshev LP, '+10 pp F1' ile 'baseline'dan kötü' arasındaki farktır.",
        "Veri artırma, attention veya kayıp mühendisliği nedeniyle çalışmaz — hepsi "
        "sabit tutulmuştur. Kazanç tamamen girdi temsilinde.",
    ]
    tb = add_bullets(slide, Cm(1.5), Cm(3), Cm(30.5), Cm(14), items, size=Pt(17))
    add_footer(slide, "Tartışma", "18 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb.shape_id])


def slide_limitations(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Sınırlamalar")
    items = [
        "Tek-veri-seti deneyi (Chapman–Shaoxing). PTB-XL veri setleri arası doğrulama "
        "önümüzdeki ayın önceliği.",
        "Altörnekleme faktörünün 500 örneğin (50 Hz) altında ki tatlı noktası "
        "henüz karakterize edilmemiştir.",
        "Yüksek frekanslı bilgi (geç-potansiyel notching, mikro-alternans) tasarım "
        "gereği kaldırılır — bunlara dayanan herhangi bir alt görev yeniden "
        "değerlendirilmelidir.",
        "Temel 1D-CNN sabittir; giriş uzunluğu ile daha derin bir CNN veya "
        "CNN + attention etkileşimi henüz ölçülmemiştir.",
        "Klinik güven skoru (tek örnek softmax) %12,9'dan %76,2'ye yükselmiştir, "
        "ancak dağıtım için sınıf-bazlı kalibrasyon hala gereklidir.",
    ]
    tb = add_bullets(slide, Cm(1.5), Cm(3), Cm(30.5), Cm(14), items, size=Pt(18))
    add_footer(slide, "Sınırlamalar", "19 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb.shape_id])


def slide_future_work(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Gelecek Çalışmalar")
    items = [
        "1.  PTB-XL üzerinde (yerli 100 Hz) altörneklemeli ve altörneklemesiz "
        "veri setleri arası doğrulama.",
        "2.  Etiket taksonomisi temizliği: ≈ 23 'EKG: X' / 'X' kopyasını birleştir "
        "(→ ~55 sınıf) ve len=500 baseline'ı yeniden çalıştır.",
        "3.  Decimate-500 girdisi üzerinde attention-CNN-LSTM tam model. Tahmin: "
        "geri kalan iyileşme çoğunlukla sınırda ST-segment / atriyal-taşikardi "
        "karışıklıklarında olacak.",
        "4.  Focal loss γ ∈ {1,2,3} grid araması ve uyarlanabilir sınıf bazında "
        "eşikleme; beklenen ekstra +1–2 pp makro F1.",
        "5.  Decimate-500 üzerinde GradCAM + SHAP: girdi zaten zamansal olarak "
        "sıkıştırıldığında attention zirveleri hâlâ P/QRS/T üzerine düşüyor mu?",
        "6.  Kenar dağıtımı: INT8 kuantizasyon + anti-aliasing'li altörnekleme, "
        "Raspberry Pi 4'e < 100 ms/örnek ile < 1 pp doğruluk kaybıyla sığmalıdır.",
        "7.  Gelecek EKG kıyaslamalarında altörneklemeyi tasarım değişkeni olarak "
        "raporla — hiper-parametre olarak değil.",
    ]
    tb = add_bullets(slide, Cm(1.5), Cm(3), Cm(30.5), Cm(14), items, size=Pt(17))
    add_footer(slide, "Gelecek çalışmalar", "20 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb.shape_id])


def slide_conclusion(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    add_title_banner(slide, "Sonuç")
    add_rect(slide, Cm(2), Cm(4), Cm(30), Cm(4.5), fill=BG)
    add_text(
        slide, Cm(3), Cm(4.8), Cm(28), Cm(3.5),
        "12 kanallı EKG çoklu etiket sınıflandırmasında, varsayılan 500 Hz × 10 s "
        "girdiyi 50 Hz × 10 s anti-aliasing'li sinyalle değiştirmek düz bir 1D-CNN'i "
        "%88,43'ten %97,34 test doğruluğuna çıkarır; çıkarım maliyeti 3,3× düşer, "
        "modelde hiçbir değişiklik olmaz.",
        size=Pt(20), italic=True, color=PRIMARY, align=PP_ALIGN.LEFT,
    )
    items = [
        "Baseline'daki en büyük tek kol, girdi temsiliydi.",
        "Altörnekleme sonrası, başlığın hibrit planı için daha büyük kol "
        "referans-düğüm veri artırmadır — kanal sayısı değil.",
        "Literatürdeki hedef rakamlar, modellemeden çok bir uzunluk-optimizasyon "
        "kalıntısı saklayabilir.",
        "Planlı tüm hibrit-model / kayıp-mühendisliği iyileştirmeleri yol haritasında "
        "kalıyor — artık çok daha güçlü bir referans noktasından başlayarak.",
    ]
    tb = add_bullets(slide, Cm(2), Cm(9.5), Cm(30), Cm(8), items, size=Pt(19))
    add_footer(slide, "Sonuç", "21 / 21")
    add_fade_transition(slide)
    add_appear_animations(slide, [tb.shape_id])


def slide_thank_you(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    sw, sh = prs.slide_width, prs.slide_height
    add_rect(slide, 0, 0, sw, sh, fill=PRIMARY)
    add_text(slide, 0, Cm(7), sw, Cm(2),
             "Teşekkürler.", size=Pt(48), bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text(slide, 0, Cm(9.5), sw, Cm(1.5),
             "Sorular için hazırım.", size=Pt(22),
             color=RGBColor(0xDC, 0xE6, 0xF5), align=PP_ALIGN.CENTER)
    add_text(slide, 0, Cm(12), sw, Cm(1),
             "Elaman Nazarkulov  ·  elaman.job@gmail.com  ·  KTMÜ",
             size=Pt(14), color=RGBColor(0xDC, 0xE6, 0xF5), align=PP_ALIGN.CENTER)
    add_fade_transition(slide)


def build_presentation() -> Path:
    prs = new_presentation()
    slide_title(prs)
    slide_motivation(prs)
    slide_research_question(prs)
    slide_related_work(prs)
    slide_dataset(prs)
    slide_method(prs)
    slide_decimation(prs)
    slide_geometry(prs)
    slide_experimental_setup(prs)
    slide_results_table(prs)
    slide_results_figure(prs)
    slide_training_histories(prs)
    slide_per_class(prs)
    slide_hybrid_intro(prs)
    slide_hybrid_augment(prs)
    slide_hybrid_channels(prs)
    slide_hybrid_deployment(prs)
    slide_speed(prs)
    slide_discussion(prs)
    slide_limitations(prs)
    slide_future_work(prs)
    slide_conclusion(prs)
    slide_thank_you(prs)
    prs.save(PPTX_PATH)
    return PPTX_PATH


if __name__ == "__main__":
    path = build_presentation()
    print("Yazıldı:", path)
